"""Office document extraction (DOCX/PPTX/XLSX).

These formats carry their own text layer, so — like a digital PDF — no
OCR engine call is needed for their native text. Embedded raster images
inside an office file are NOT OCR'd in this version (see README
limitations); only the document's own text/table content is extracted.

Security: python-docx/python-pptx/openpyxl parse the OOXML (zip+XML)
structure only; no macros are executed and no external file is opened
by these libraries, satisfying "never execute uploaded files".
"""

from __future__ import annotations

from pathlib import Path

from Agents.ocr_agent.exceptions import DocumentReadError, PageExtractionError


def extract_docx_pages(path: Path) -> list[str]:
    try:
        import docx  # python-docx
    except ImportError as exc:
        raise DocumentReadError(
            "python-docx is required to process .docx files. Install it with "
            "`pip install python-docx`."
        ) from exc

    try:
        document = docx.Document(str(path))
    except Exception as exc:
        raise DocumentReadError(f"Failed to open DOCX: {exc}") from exc

    parts: list[str] = []
    for para in document.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))

    # DOCX has no native "page" boundary in the XML; treat the whole
    # document as a single logical page (still processed in full).
    text = "\n".join(parts)
    return [text] if text.strip() else [""]


def extract_pptx_pages(path: Path) -> list[str]:
    try:
        from pptx import Presentation  # python-pptx
    except ImportError as exc:
        raise DocumentReadError(
            "python-pptx is required to process .pptx files. Install it with "
            "`pip install python-pptx`."
        ) from exc

    try:
        presentation = Presentation(str(path))
    except Exception as exc:
        raise DocumentReadError(f"Failed to open PPTX: {exc}") from exc

    pages: list[str] = []
    for slide in presentation.slides:
        lines: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame is not None:
                text = shape.text_frame.text
                if text and text.strip():
                    lines.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        lines.append(" | ".join(cells))
        pages.append("\n".join(lines))
    if not pages:
        raise PageExtractionError("PPTX contains no slides.")
    return pages


def extract_xlsx_pages(path: Path) -> list[str]:
    try:
        import openpyxl
    except ImportError as exc:
        raise DocumentReadError(
            "openpyxl is required to process .xlsx files. Install it with "
            "`pip install openpyxl`."
        ) from exc

    try:
        workbook = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    except Exception as exc:
        raise DocumentReadError(f"Failed to open XLSX: {exc}") from exc

    pages: list[str] = []
    try:
        for sheet in workbook.worksheets:
            rows: list[str] = [f"# {sheet.title}"]
            for row in sheet.iter_rows(values_only=True):
                values = ["" if v is None else str(v) for v in row]
                if any(v.strip() for v in values):
                    rows.append(" | ".join(values))
            pages.append("\n".join(rows))
    finally:
        workbook.close()
    if not pages:
        raise PageExtractionError("XLSX contains no sheets.")
    return pages
